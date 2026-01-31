"""PowerPoint translator module - JA↔VI translation preserving formatting.

Uses python-pptx for cross-platform support (no PowerPoint installation required).

Usage:
    python scripts/pptx_translator.py input.pptx --to ja
    python scripts/pptx_translator.py input.pptx --to vi --glossary glossaries/custom.json
"""

import os
import sys
import time
import re
import argparse
import logging
from pathlib import Path
from datetime import datetime
from typing import List, Tuple, Optional, Dict, Any

try:
    from pptx import Presentation
    from pptx.util import Pt
    from openai import OpenAI
    from dotenv import load_dotenv
except ImportError as e:
    print(f"❌ Missing dependency: {e}")
    print("Install with: pip install python-pptx openai python-dotenv")
    sys.exit(1)

# Load environment
load_dotenv()

# Constants
API_DELAY = 2  # seconds between API calls
BATCH_SIZE = 30  # max texts per batch (smaller than Excel due to longer content)


class PptxTranslator:
    """Translate PowerPoint files while preserving formatting."""

    def __init__(
        self,
        api_key: Optional[str] = None,
        glossary: Optional[Dict[str, str]] = None,
        system_prompt: Optional[str] = None,
        log_dir: Optional[str] = None
    ):
        """Initialize translator.

        Args:
            api_key: Gemini API key (or from GEMINI_API_KEY env)
            glossary: Dictionary of term translations
            system_prompt: Custom system prompt for translation
            log_dir: Directory for log files (default: no logging)
        """
        self.api_key = api_key or os.getenv("GEMINI_API_KEY")
        if not self.api_key:
            raise ValueError("GEMINI_API_KEY environment variable required")

        self.client = OpenAI(
            base_url="https://generativelanguage.googleapis.com/v1beta/openai/",
            api_key=self.api_key,
        )
        self.glossary = glossary or {}
        self.system_prompt = system_prompt or self._load_default_prompt()
        self.logger = self._setup_logging(log_dir) if log_dir else None

    def _setup_logging(self, log_dir: str) -> logging.Logger:
        """Setup logging to file."""
        os.makedirs(log_dir, exist_ok=True)
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        log_file = os.path.join(log_dir, f'pptx_translation_{timestamp}.log')

        logger = logging.getLogger('pptx_translator')
        logger.setLevel(logging.INFO)
        handler = logging.FileHandler(log_file, encoding='utf-8')
        handler.setFormatter(logging.Formatter('%(asctime)s - %(levelname)s - %(message)s'))
        logger.addHandler(handler)
        return logger

    def _log(self, level: str, message: str):
        """Log message if logger is configured."""
        if self.logger:
            getattr(self.logger, level)(message)

    def _load_default_prompt(self) -> str:
        """Load default system prompt from file."""
        prompt_file = Path(__file__).parent / "pptx_system_prompt.txt"
        if prompt_file.exists():
            return prompt_file.read_text(encoding='utf-8')
        return self._get_fallback_prompt()

    def _get_fallback_prompt(self) -> str:
        """Fallback prompt if file not found."""
        return """You are a professional translator for IT projects.
Translate accurately while preserving formatting (bullets, line breaks, spacing).
For words with underscores (ABC_Order_Management), treat as separate words.
Output ONLY the translated text, no explanations.
Separate each translation with "---" on its own line."""

    def _build_user_prompt(self, texts: List[str], target_lang: str) -> str:
        """Build user prompt with glossary context."""
        direction = "Vietnamese to Japanese" if target_lang == "ja" else "Japanese to Vietnamese"
        separator = "\n---\n"
        combined = separator.join(texts)

        # Add glossary context if available
        glossary_note = ""
        if self.glossary:
            terms = "\n".join(f"- {k} → {v}" for k, v in list(self.glossary.items())[:30])
            glossary_note = f"\n\nGLOSSARY (use these exact translations):\n{terms}\n"

        return f"""Translate from {direction}. Preserve all formatting.
{glossary_note}
Output each translation separated by "---" on its own line.

TEXTS:
{combined}"""

    def translate_batch(self, texts: List[str], target_lang: str = "ja") -> List[str]:
        """Translate a batch of texts.

        Args:
            texts: List of texts to translate
            target_lang: Target language ('ja' or 'vi')

        Returns:
            List of translated texts
        """
        if not texts:
            return []

        self._log('info', f"Translating batch of {len(texts)} texts")

        try:
            response = self.client.chat.completions.create(
                model="gemini-2.5-flash-lite",
                messages=[
                    {"role": "system", "content": self.system_prompt},
                    {"role": "user", "content": self._build_user_prompt(texts, target_lang)}
                ],
                temperature=0.3
            )

            translated = response.choices[0].message.content.strip()
            parts = translated.split("\n---\n")

            # Handle count mismatch
            if len(parts) != len(texts):
                self._log('warning', f"Translation count mismatch: {len(parts)} vs {len(texts)}")
                if len(parts) < len(texts):
                    parts.extend([""] * (len(texts) - len(parts)))
                else:
                    parts = parts[:len(texts)]

            time.sleep(API_DELAY)
            return parts

        except Exception as e:
            self._log('error', f"Translation error: {e}")
            print(f"❌ Translation error: {e}")
            return texts  # fallback to original

    @staticmethod
    def _split_text_by_paragraphs(text: str) -> List[str]:
        """Split text into paragraphs, handling bullet points."""
        lines = text.split('\n')
        result = []
        current_text = ""

        for line in lines:
            line = line.strip()
            if not line:
                if current_text:
                    result.append(current_text)
                    current_text = ""
                continue

            # Check if line starts with bullet or numbering
            if re.match(r'^[•\-\*]|\d+[.)]', line):
                if current_text:
                    result.append(current_text)
                current_text = line
            elif current_text:
                if re.match(r'^[•\-\*]|\d+[.)]', current_text.split('\n')[0]):
                    current_text += '\n' + line
                else:
                    result.append(current_text)
                    current_text = line
            else:
                current_text = line

        if current_text:
            result.append(current_text)

        return result

    def _extract_table_texts(self, shape: Any) -> Tuple[List[str], List[Tuple]]:
        """Extract texts from a table shape."""
        texts = []
        locations = []

        if not hasattr(shape, "table"):
            return texts, locations

        for row_idx, row in enumerate(shape.table.rows):
            for cell_idx, cell in enumerate(row.cells):
                if cell.text.strip():
                    para_texts = self._split_text_by_paragraphs(cell.text)
                    for i, para_text in enumerate(para_texts):
                        texts.append(para_text)
                        locations.append((row_idx, cell_idx, i))

        return texts, locations

    def translate_file(
        self,
        input_path: str,
        output_path: Optional[str] = None,
        target_lang: str = "ja",
        target_font: str = "Meiryo UI"
    ) -> Optional[str]:
        """Translate PowerPoint file.

        Args:
            input_path: Path to input PPTX file
            output_path: Path to output file (auto-generated if None)
            target_lang: Target language ('ja' or 'vi')
            target_font: Font for translated text (default: Meiryo UI for Japanese)

        Returns:
            Output file path or None on error
        """
        input_path = Path(input_path)
        if not input_path.exists():
            print(f"❌ File not found: {input_path}")
            return None

        # Generate output path
        if not output_path:
            output_dir = input_path.parent / "output"
            output_dir.mkdir(exist_ok=True)
            lang_suffix = "_ja" if target_lang == "ja" else "_vi"
            output_path = output_dir / f"{input_path.stem}{lang_suffix}.pptx"
        output_path = Path(output_path)

        # Set font based on target language
        if target_lang == "vi":
            target_font = "Arial"  # Vietnamese-friendly font

        print(f"\n🔄 Processing: {input_path.name}")
        print(f"🎯 Target: {'Japanese' if target_lang == 'ja' else 'Vietnamese'}")

        try:
            prs = Presentation(str(input_path))
            all_texts = []
            text_locations = []

            # Extract all texts
            for slide_idx, slide in enumerate(prs.slides):
                for shape_idx, shape in enumerate(slide.shapes):
                    # Handle regular text shapes
                    if hasattr(shape, "text") and shape.text.strip():
                        if hasattr(shape, "text_frame"):
                            for para_idx, para in enumerate(shape.text_frame.paragraphs):
                                if para.text.strip():
                                    all_texts.append(para.text.strip())
                                    text_locations.append(("paragraph", slide_idx, shape_idx, para_idx))
                        else:
                            split_texts = self._split_text_by_paragraphs(shape.text)
                            for para_text in split_texts:
                                all_texts.append(para_text)
                                text_locations.append(("text", slide_idx, shape_idx, len(all_texts)))

                    # Handle tables
                    if hasattr(shape, "table"):
                        table_texts, table_locs = self._extract_table_texts(shape)
                        all_texts.extend(table_texts)
                        for (row_idx, cell_idx, para_idx) in table_locs:
                            text_locations.append(("table", slide_idx, shape_idx, row_idx, cell_idx, para_idx))

            if not all_texts:
                print(f"   ⚠️ No text found")
                return None

            # Translate in batches
            translated_texts = []
            batches = [all_texts[i:i+BATCH_SIZE] for i in range(0, len(all_texts), BATCH_SIZE)]

            print(f"   📦 {len(all_texts)} texts in {len(batches)} batches")

            for i, batch in enumerate(batches):
                progress = (i + 1) / len(batches) * 100
                print(f"   🔄 Batch {i+1}/{len(batches)} [{int(progress)}%]")
                translations = self.translate_batch(batch, target_lang)
                translated_texts.extend(translations)

            # Update presentation
            for location, translated_text in zip(text_locations, translated_texts):
                if not translated_text:
                    continue

                if location[0] == "paragraph":
                    _, slide_idx, shape_idx, para_idx = location
                    shape = prs.slides[slide_idx].shapes[shape_idx]
                    if hasattr(shape, "text_frame") and para_idx < len(shape.text_frame.paragraphs):
                        para = shape.text_frame.paragraphs[para_idx]

                        # Store original formatting
                        orig_align = para.alignment
                        orig_level = para.level
                        orig_font_sizes = [run.font.size for run in para.runs if hasattr(run.font, 'size')]

                        para.text = translated_text

                        # Apply font and restore formatting
                        for idx, run in enumerate(para.runs):
                            run.font.name = target_font
                            if idx < len(orig_font_sizes) and orig_font_sizes[idx]:
                                run.font.size = orig_font_sizes[idx]

                        para.alignment = orig_align
                        para.level = orig_level

                elif location[0] == "table":
                    _, slide_idx, shape_idx, row_idx, cell_idx, para_idx = location
                    shape = prs.slides[slide_idx].shapes[shape_idx]
                    if hasattr(shape, "table"):
                        cell = shape.table.rows[row_idx].cells[cell_idx]
                        if hasattr(cell, "text_frame") and para_idx < len(cell.text_frame.paragraphs):
                            para = cell.text_frame.paragraphs[para_idx]

                            orig_align = para.alignment
                            orig_level = para.level
                            orig_font_sizes = [run.font.size for run in para.runs if hasattr(run.font, 'size')]

                            para.text = translated_text

                            for idx, run in enumerate(para.runs):
                                run.font.name = target_font
                                if idx < len(orig_font_sizes) and orig_font_sizes[idx]:
                                    run.font.size = orig_font_sizes[idx]

                            para.alignment = orig_align
                            para.level = orig_level

            # Save
            print(f"\n💾 Saving: {output_path}")
            prs.save(str(output_path))
            print(f"✅ Done: {output_path}")
            return str(output_path)

        except Exception as e:
            self._log('error', f"Error: {e}")
            print(f"❌ Error: {e}")
            return None


def main():
    """CLI entry point."""
    parser = argparse.ArgumentParser(description='Translate PowerPoint files JA↔VI')
    parser.add_argument('input', help='Input PPTX file path')
    parser.add_argument('--to', choices=['ja', 'vi'], default='ja',
                        help='Target language (default: ja)')
    parser.add_argument('--output', '-o', help='Output file path')
    parser.add_argument('--glossary', '-g', help='Glossary JSON file path')
    parser.add_argument('--log-dir', help='Directory for log files')

    args = parser.parse_args()

    # Load glossary if provided
    glossary = {}
    if args.glossary:
        import json
        glossary_path = Path(args.glossary)
        if glossary_path.exists():
            glossary = json.loads(glossary_path.read_text(encoding='utf-8'))
            print(f"📖 Loaded {len(glossary)} glossary terms")

    # Translate
    translator = PptxTranslator(glossary=glossary, log_dir=args.log_dir)
    result = translator.translate_file(args.input, args.output, args.to)

    if result:
        print(f"\n✅ Translation complete: {result}")
    else:
        print("\n❌ Translation failed")
        sys.exit(1)


if __name__ == '__main__':
    main()
